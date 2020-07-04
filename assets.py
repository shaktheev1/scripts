from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
import csv
import os
from wand.image import Image

images='images'
directory1=input("Enter folder name:")
directory='/mnt/art/upload/PPT/'+directory1


if not(os.path.exists(directory)):
    print("'{}' does not exist. Please check.".format(directory))
    exit()

paths=Path(directory).glob('**/*.pptx')

x=0
n=0
total=0
total_slides=0
filepath=''
chapter_number=''
images_sig={}
slide_sig={}
no_of_slides=0
slide_for_dup=0
duplicate_file = os.path.join(directory,'duplicates.txt')
duplicates=open(duplicate_file,"w") 

def find_duplicate_image(image):
    img=Image(filename=image)
    sig=img.signature
    if sig in images_sig.values():
        images_sig[image]=sig
        slide_sig[image]=slide_for_dup
        duplicates.write("{} on slide {} is a duplicate of following image: \n".format(image, slide_sig[image]))
        found=list(images_sig.keys())[list(images_sig.values()).index(sig)]
        duplicates.write("{} on slide {}".format(found,slide_sig[found]))
        duplicates.write("\n\n")
    else:
        images_sig[image]=sig
        slide_sig[image]=slide_for_dup
    
def write_image(shape):
    global n
    global filepath

    # Creating 'images' folder
    images1=os.path.join(directory, 'images')
    if not(os.path.exists(images1)):
        os.mkdir(images1)
        os.chmod(os.path.join(directory,images1), 0o02775)
        os.system("chgrp -R space "+os.path.join(directory,images1))

    # Creating 'chapter' folder
    images=os.path.join(images1, chapter_number)
    if not(os.path.exists(images)):
        os.mkdir(images)
        os.chmod(os.path.join(directory,images), 0o02775)
        os.system("chgrp -R space "+os.path.join(directory,images))

    image = shape.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    # Joining input folder and image folders

    image_filename = os.path.join(directory,'images/{}/image{:03d}.{}'.format(chapter_number, n, image.ext))

    n += 1

# print(image_filename)
    
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)

    find_duplicate_image(image_filename)

def visitor(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        write_image(shape)

def iter_picture_shapes(prs):
    global slide_for_dup
    slide_for_dup=0
    for slide in prs.slides:
        slide_for_dup=slide_for_dup+1
        for shape in slide.shapes:
            visitor(shape)

# Joining input folder and image_count.csv
csv_path = os.path.join(directory,'image_count.csv')

with open(csv_path, 'w') as csvfile:
    csvwriter = csv.writer(csvfile)
    #Adding headet to CSV
    csvwriter.writerow(['Chapter Number', 'Number of Images', 'Number of Slides'])            
    for filepath in paths:
        # Splitting the filepath
        mystr=os.path.basename(filepath)
        chapter_number=os.path.splitext(mystr)[0]
        # Slide count
        iter_picture_shapes(Presentation(filepath))
        no_of_slides=len(Presentation(filepath).slides)
        # Writing each row to CSV
        csvwriter.writerow([chapter_number, n, no_of_slides])
        # Aggregating number of images in each chapter
        total = total + n
        # Aggregating number of slides in each presentation
        total_slides = total_slides + no_of_slides
        # Resetting image count to zero for each presentation
        n=0
    # Appending total to CSV    
    csvwriter.writerow(['Total',total,total_slides])

duplicates.close()
print("Successfull! 'image_count.csv' and 'duplicates.txt' are stored inside {}".format(directory))


# result = get_details()
# print(result)
